"""Routes pour les analyses LLM (résumé, points clés, etc.)."""

import io
import json
import time

from fastapi import APIRouter, HTTPException, Depends, Query, Request
from fastapi.responses import StreamingResponse, Response

from backend import config
from backend.api.deps import verify_api_key
from backend.database import (
    get_connection, get_transcription_with_segments,
    get_analysis, get_analyses, save_analysis_sync,
)
from backend.llm_processing import ollama_client
from backend.llm_processing.summarizer import summarize_stream
from backend.llm_processing.key_points import extract_key_points_stream
from backend.llm_processing.actions import extract_actions_stream
from backend.llm_processing.study_cards import generate_study_cards_stream
from backend.llm_processing.quiz import generate_quiz_stream
from backend.llm_processing.mindmap import generate_mindmap_stream
from backend.llm_processing.slides import generate_slides_stream
from backend.llm_processing.infographic import generate_infographic_stream
from backend.llm_processing.data_table import extract_data_table_stream

router = APIRouter()


async def _parse_instructions(request: Request) -> str | None:
    """Parse optional custom instructions from JSON body."""
    try:
        body = await request.json()
        return body.get("instructions") if isinstance(body, dict) else None
    except Exception:
        return None


@router.get("/api/llm/status")
async def llm_status():
    """Vérifie si le LLM (Ollama) est disponible."""
    available = ollama_client.is_available()
    models = ollama_client.list_models() if available else []
    return {
        "available": available,
        "models": models,
        "current_model": config.LLM_MODEL,
    }


@router.post("/api/transcriptions/{tid}/summarize", dependencies=[Depends(verify_api_key)])
async def summarize(tid: int, request: Request, model: str = Query(None)):
    """Génère un résumé via LLM en streaming SSE."""
    instructions = await _parse_instructions(request)
    if not ollama_client.is_available():
        raise HTTPException(503, "LLM non disponible. Vérifiez qu'Ollama est lancé.")

    db = await get_connection(config.DB_PATH)
    try:
        t = await get_transcription_with_segments(db, tid)
        if not t:
            raise HTTPException(404, "Transcription non trouvée.")
    finally:
        await db.close()

    segments = t.get("segments", [])
    if not segments:
        raise HTTPException(400, "Aucun segment à résumer.")

    full_text = " ".join(s["text"] for s in segments if s.get("text"))
    if not full_text.strip():
        raise HTTPException(400, "Le texte de la transcription est vide.")

    filename = t.get("filename", "audio")
    use_model = model or config.LLM_MODEL

    def generate():
        start_time = time.time()
        try:
            yield f"data: {json.dumps({'type': 'status', 'message': 'Génération du résumé...'}, ensure_ascii=False)}\n\n"

            for chunk in summarize_stream(full_text, filename=filename, model=use_model, custom_instructions=instructions):
                if chunk["done"]:
                    # Sauvegarder en DB
                    processing_ms = int((time.time() - start_time) * 1000)
                    analysis_id = save_analysis_sync(
                        config.DB_PATH, tid, "summary",
                        chunk["full_text"], use_model, processing_ms,
                    )
                    yield f"data: {json.dumps({'type': 'done', 'analysis_id': analysis_id, 'processing_ms': processing_ms}, ensure_ascii=False)}\n\n"
                else:
                    yield f"data: {json.dumps({'type': 'token', 'token': chunk['token']}, ensure_ascii=False)}\n\n"

        except Exception as e:
            yield f"data: {json.dumps({'type': 'error', 'message': str(e)}, ensure_ascii=False)}\n\n"

    return StreamingResponse(generate(), media_type="text/event-stream")


@router.post("/api/transcriptions/{tid}/extract-key-points", dependencies=[Depends(verify_api_key)])
async def extract_key_points(tid: int, request: Request, model: str = Query(None)):
    """Extrait les points cles via LLM en streaming SSE."""
    instructions = await _parse_instructions(request)
    if not ollama_client.is_available():
        raise HTTPException(503, "LLM non disponible. Verifiez qu'Ollama est lance.")

    db = await get_connection(config.DB_PATH)
    try:
        t = await get_transcription_with_segments(db, tid)
        if not t:
            raise HTTPException(404, "Transcription non trouvee.")
    finally:
        await db.close()

    segments = t.get("segments", [])
    if not segments:
        raise HTTPException(400, "Aucun segment a analyser.")

    full_text = " ".join(s["text"] for s in segments if s.get("text"))
    if not full_text.strip():
        raise HTTPException(400, "Le texte de la transcription est vide.")

    filename = t.get("filename", "audio")
    use_model = model or config.LLM_MODEL

    def generate():
        start_time = time.time()
        try:
            yield f"data: {json.dumps({'type': 'status', 'message': 'Extraction des points cles...'}, ensure_ascii=False)}\n\n"

            for chunk in extract_key_points_stream(full_text, filename=filename, model=use_model, custom_instructions=instructions):
                if chunk["done"]:
                    processing_ms = int((time.time() - start_time) * 1000)
                    analysis_id = save_analysis_sync(
                        config.DB_PATH, tid, "key_points",
                        chunk["full_text"], use_model, processing_ms,
                    )
                    yield f"data: {json.dumps({'type': 'done', 'analysis_id': analysis_id, 'processing_ms': processing_ms}, ensure_ascii=False)}\n\n"
                else:
                    yield f"data: {json.dumps({'type': 'token', 'token': chunk['token']}, ensure_ascii=False)}\n\n"

        except Exception as e:
            yield f"data: {json.dumps({'type': 'error', 'message': str(e)}, ensure_ascii=False)}\n\n"

    return StreamingResponse(generate(), media_type="text/event-stream")


@router.post("/api/transcriptions/{tid}/extract-actions", dependencies=[Depends(verify_api_key)])
async def extract_actions(tid: int, request: Request, model: str = Query(None)):
    """Extrait les actions via LLM en streaming SSE."""
    instructions = await _parse_instructions(request)
    if not ollama_client.is_available():
        raise HTTPException(503, "LLM non disponible. Verifiez qu'Ollama est lance.")

    db = await get_connection(config.DB_PATH)
    try:
        t = await get_transcription_with_segments(db, tid)
        if not t:
            raise HTTPException(404, "Transcription non trouvee.")
    finally:
        await db.close()

    segments = t.get("segments", [])
    if not segments:
        raise HTTPException(400, "Aucun segment a analyser.")

    full_text = " ".join(s["text"] for s in segments if s.get("text"))
    if not full_text.strip():
        raise HTTPException(400, "Le texte de la transcription est vide.")

    filename = t.get("filename", "audio")
    use_model = model or config.LLM_MODEL

    def generate():
        start_time = time.time()
        try:
            yield f"data: {json.dumps({'type': 'status', 'message': 'Extraction des actions...'}, ensure_ascii=False)}\n\n"

            for chunk in extract_actions_stream(full_text, filename=filename, model=use_model, custom_instructions=instructions):
                if chunk["done"]:
                    processing_ms = int((time.time() - start_time) * 1000)
                    analysis_id = save_analysis_sync(
                        config.DB_PATH, tid, "actions",
                        chunk["full_text"], use_model, processing_ms,
                    )
                    yield f"data: {json.dumps({'type': 'done', 'analysis_id': analysis_id, 'processing_ms': processing_ms}, ensure_ascii=False)}\n\n"
                else:
                    yield f"data: {json.dumps({'type': 'token', 'token': chunk['token']}, ensure_ascii=False)}\n\n"

        except Exception as e:
            yield f"data: {json.dumps({'type': 'error', 'message': str(e)}, ensure_ascii=False)}\n\n"

    return StreamingResponse(generate(), media_type="text/event-stream")


@router.post("/api/transcriptions/{tid}/generate-study-cards", dependencies=[Depends(verify_api_key)])
async def study_cards(tid: int, request: Request, model: str = Query(None)):
    """Genere des fiches d'apprentissage via LLM en streaming SSE."""
    instructions = await _parse_instructions(request)
    if not ollama_client.is_available():
        raise HTTPException(503, "LLM non disponible. Verifiez qu'Ollama est lance.")

    db = await get_connection(config.DB_PATH)
    try:
        t = await get_transcription_with_segments(db, tid)
        if not t:
            raise HTTPException(404, "Transcription non trouvee.")
    finally:
        await db.close()

    segments = t.get("segments", [])
    if not segments:
        raise HTTPException(400, "Aucun segment a analyser.")

    full_text = " ".join(s["text"] for s in segments if s.get("text"))
    if not full_text.strip():
        raise HTTPException(400, "Le texte de la transcription est vide.")

    filename = t.get("filename", "audio")
    use_model = model or config.LLM_MODEL

    def generate():
        start_time = time.time()
        try:
            yield f"data: {json.dumps({'type': 'status', 'message': 'Generation des fiches...'}, ensure_ascii=False)}\n\n"

            for chunk in generate_study_cards_stream(full_text, filename=filename, model=use_model, custom_instructions=instructions):
                if chunk["done"]:
                    processing_ms = int((time.time() - start_time) * 1000)
                    analysis_id = save_analysis_sync(
                        config.DB_PATH, tid, "study_cards",
                        chunk["full_text"], use_model, processing_ms,
                    )
                    yield f"data: {json.dumps({'type': 'done', 'analysis_id': analysis_id, 'processing_ms': processing_ms}, ensure_ascii=False)}\n\n"
                else:
                    yield f"data: {json.dumps({'type': 'token', 'token': chunk['token']}, ensure_ascii=False)}\n\n"

        except Exception as e:
            yield f"data: {json.dumps({'type': 'error', 'message': str(e)}, ensure_ascii=False)}\n\n"

    return StreamingResponse(generate(), media_type="text/event-stream")


@router.post("/api/transcriptions/{tid}/generate-quiz", dependencies=[Depends(verify_api_key)])
async def quiz(tid: int, request: Request, model: str = Query(None)):
    """Genere un quiz via LLM en streaming SSE."""
    instructions = await _parse_instructions(request)
    if not ollama_client.is_available():
        raise HTTPException(503, "LLM non disponible. Verifiez qu'Ollama est lance.")

    db = await get_connection(config.DB_PATH)
    try:
        t = await get_transcription_with_segments(db, tid)
        if not t:
            raise HTTPException(404, "Transcription non trouvee.")
    finally:
        await db.close()

    segments = t.get("segments", [])
    if not segments:
        raise HTTPException(400, "Aucun segment a analyser.")

    full_text = " ".join(s["text"] for s in segments if s.get("text"))
    if not full_text.strip():
        raise HTTPException(400, "Le texte de la transcription est vide.")

    filename = t.get("filename", "audio")
    use_model = model or config.LLM_MODEL

    def generate():
        start_time = time.time()
        try:
            yield f"data: {json.dumps({'type': 'status', 'message': 'Generation du quiz...'}, ensure_ascii=False)}\n\n"

            for chunk in generate_quiz_stream(full_text, filename=filename, model=use_model, custom_instructions=instructions):
                if chunk["done"]:
                    processing_ms = int((time.time() - start_time) * 1000)
                    analysis_id = save_analysis_sync(
                        config.DB_PATH, tid, "quiz",
                        chunk["full_text"], use_model, processing_ms,
                    )
                    yield f"data: {json.dumps({'type': 'done', 'analysis_id': analysis_id, 'processing_ms': processing_ms}, ensure_ascii=False)}\n\n"
                else:
                    yield f"data: {json.dumps({'type': 'token', 'token': chunk['token']}, ensure_ascii=False)}\n\n"

        except Exception as e:
            yield f"data: {json.dumps({'type': 'error', 'message': str(e)}, ensure_ascii=False)}\n\n"

    return StreamingResponse(generate(), media_type="text/event-stream")


@router.post("/api/transcriptions/{tid}/generate-mindmap", dependencies=[Depends(verify_api_key)])
async def mindmap(tid: int, request: Request, model: str = Query(None)):
    """Genere une carte mentale via LLM en streaming SSE."""
    instructions = await _parse_instructions(request)
    if not ollama_client.is_available():
        raise HTTPException(503, "LLM non disponible. Verifiez qu'Ollama est lance.")

    db = await get_connection(config.DB_PATH)
    try:
        t = await get_transcription_with_segments(db, tid)
        if not t:
            raise HTTPException(404, "Transcription non trouvee.")
    finally:
        await db.close()

    segments = t.get("segments", [])
    if not segments:
        raise HTTPException(400, "Aucun segment a analyser.")

    full_text = " ".join(s["text"] for s in segments if s.get("text"))
    if not full_text.strip():
        raise HTTPException(400, "Le texte de la transcription est vide.")

    filename = t.get("filename", "audio")
    use_model = model or config.LLM_MODEL

    def generate():
        start_time = time.time()
        try:
            yield f"data: {json.dumps({'type': 'status', 'message': 'Generation de la carte mentale...'}, ensure_ascii=False)}\n\n"

            for chunk in generate_mindmap_stream(full_text, filename=filename, model=use_model, custom_instructions=instructions):
                if chunk["done"]:
                    processing_ms = int((time.time() - start_time) * 1000)
                    analysis_id = save_analysis_sync(
                        config.DB_PATH, tid, "mindmap",
                        chunk["full_text"], use_model, processing_ms,
                    )
                    yield f"data: {json.dumps({'type': 'done', 'analysis_id': analysis_id, 'processing_ms': processing_ms}, ensure_ascii=False)}\n\n"
                else:
                    yield f"data: {json.dumps({'type': 'token', 'token': chunk['token']}, ensure_ascii=False)}\n\n"

        except Exception as e:
            yield f"data: {json.dumps({'type': 'error', 'message': str(e)}, ensure_ascii=False)}\n\n"

    return StreamingResponse(generate(), media_type="text/event-stream")


@router.post("/api/transcriptions/{tid}/generate-slides", dependencies=[Depends(verify_api_key)])
async def slides(tid: int, request: Request, model: str = Query(None)):
    """Genere des slides via LLM en streaming SSE."""
    instructions = await _parse_instructions(request)
    if not ollama_client.is_available():
        raise HTTPException(503, "LLM non disponible. Verifiez qu'Ollama est lance.")

    db = await get_connection(config.DB_PATH)
    try:
        t = await get_transcription_with_segments(db, tid)
        if not t:
            raise HTTPException(404, "Transcription non trouvee.")
    finally:
        await db.close()

    segments = t.get("segments", [])
    if not segments:
        raise HTTPException(400, "Aucun segment a analyser.")

    full_text = " ".join(s["text"] for s in segments if s.get("text"))
    if not full_text.strip():
        raise HTTPException(400, "Le texte de la transcription est vide.")

    filename = t.get("filename", "audio")
    use_model = model or config.LLM_MODEL

    def generate():
        start_time = time.time()
        try:
            yield f"data: {json.dumps({'type': 'status', 'message': 'Generation des slides...'}, ensure_ascii=False)}\n\n"

            for chunk in generate_slides_stream(full_text, filename=filename, model=use_model, custom_instructions=instructions):
                if chunk["done"]:
                    processing_ms = int((time.time() - start_time) * 1000)
                    analysis_id = save_analysis_sync(
                        config.DB_PATH, tid, "slides",
                        chunk["full_text"], use_model, processing_ms,
                    )
                    yield f"data: {json.dumps({'type': 'done', 'analysis_id': analysis_id, 'processing_ms': processing_ms}, ensure_ascii=False)}\n\n"
                else:
                    yield f"data: {json.dumps({'type': 'token', 'token': chunk['token']}, ensure_ascii=False)}\n\n"

        except Exception as e:
            yield f"data: {json.dumps({'type': 'error', 'message': str(e)}, ensure_ascii=False)}\n\n"

    return StreamingResponse(generate(), media_type="text/event-stream")


@router.post("/api/transcriptions/{tid}/generate-infographic", dependencies=[Depends(verify_api_key)])
async def infographic(tid: int, request: Request, model: str = Query(None)):
    """Genere une infographie via LLM en streaming SSE."""
    instructions = await _parse_instructions(request)
    if not ollama_client.is_available():
        raise HTTPException(503, "LLM non disponible. Verifiez qu'Ollama est lance.")

    db = await get_connection(config.DB_PATH)
    try:
        t = await get_transcription_with_segments(db, tid)
        if not t:
            raise HTTPException(404, "Transcription non trouvee.")
    finally:
        await db.close()

    segments = t.get("segments", [])
    if not segments:
        raise HTTPException(400, "Aucun segment a analyser.")

    full_text = " ".join(s["text"] for s in segments if s.get("text"))
    if not full_text.strip():
        raise HTTPException(400, "Le texte de la transcription est vide.")

    filename = t.get("filename", "audio")
    use_model = model or config.LLM_MODEL

    def generate():
        start_time = time.time()
        try:
            yield f"data: {json.dumps({'type': 'status', 'message': 'Generation de l infographie...'}, ensure_ascii=False)}\n\n"

            for chunk in generate_infographic_stream(full_text, filename=filename, model=use_model, custom_instructions=instructions):
                if chunk["done"]:
                    processing_ms = int((time.time() - start_time) * 1000)
                    analysis_id = save_analysis_sync(
                        config.DB_PATH, tid, "infographic",
                        chunk["full_text"], use_model, processing_ms,
                    )
                    yield f"data: {json.dumps({'type': 'done', 'analysis_id': analysis_id, 'processing_ms': processing_ms}, ensure_ascii=False)}\n\n"
                else:
                    yield f"data: {json.dumps({'type': 'token', 'token': chunk['token']}, ensure_ascii=False)}\n\n"

        except Exception as e:
            yield f"data: {json.dumps({'type': 'error', 'message': str(e)}, ensure_ascii=False)}\n\n"

    return StreamingResponse(generate(), media_type="text/event-stream")


@router.post("/api/transcriptions/{tid}/extract-data-table", dependencies=[Depends(verify_api_key)])
async def data_table(tid: int, request: Request, model: str = Query(None)):
    """Extrait des tableaux de donnees via LLM en streaming SSE."""
    instructions = await _parse_instructions(request)
    if not ollama_client.is_available():
        raise HTTPException(503, "LLM non disponible. Verifiez qu'Ollama est lance.")

    db = await get_connection(config.DB_PATH)
    try:
        t = await get_transcription_with_segments(db, tid)
        if not t:
            raise HTTPException(404, "Transcription non trouvee.")
    finally:
        await db.close()

    segments = t.get("segments", [])
    if not segments:
        raise HTTPException(400, "Aucun segment a analyser.")

    full_text = " ".join(s["text"] for s in segments if s.get("text"))
    if not full_text.strip():
        raise HTTPException(400, "Le texte de la transcription est vide.")

    filename = t.get("filename", "audio")
    use_model = model or config.LLM_MODEL

    def generate():
        start_time = time.time()
        try:
            yield f"data: {json.dumps({'type': 'status', 'message': 'Extraction des tableaux de donnees...'}, ensure_ascii=False)}\n\n"

            for chunk in extract_data_table_stream(full_text, filename=filename, model=use_model, custom_instructions=instructions):
                if chunk["done"]:
                    processing_ms = int((time.time() - start_time) * 1000)
                    analysis_id = save_analysis_sync(
                        config.DB_PATH, tid, "data_table",
                        chunk["full_text"], use_model, processing_ms,
                    )
                    yield f"data: {json.dumps({'type': 'done', 'analysis_id': analysis_id, 'processing_ms': processing_ms}, ensure_ascii=False)}\n\n"
                else:
                    yield f"data: {json.dumps({'type': 'token', 'token': chunk['token']}, ensure_ascii=False)}\n\n"

        except Exception as e:
            yield f"data: {json.dumps({'type': 'error', 'message': str(e)}, ensure_ascii=False)}\n\n"

    return StreamingResponse(generate(), media_type="text/event-stream")


@router.get("/api/transcriptions/{tid}/analyses", dependencies=[Depends(verify_api_key)])
async def list_analyses(tid: int):
    """Liste toutes les analyses d'une transcription."""
    db = await get_connection(config.DB_PATH)
    try:
        analyses = await get_analyses(db, tid)
        return {"analyses": analyses}
    finally:
        await db.close()


@router.get("/api/transcriptions/{tid}/analysis/{analysis_type}", dependencies=[Depends(verify_api_key)])
async def get_analysis_by_type(tid: int, analysis_type: str):
    """Récupère la dernière analyse d'un type donné."""
    db = await get_connection(config.DB_PATH)
    try:
        analysis = await get_analysis(db, tid, analysis_type)
        if not analysis:
            raise HTTPException(404, f"Aucune analyse de type '{analysis_type}' trouvée.")
        return analysis
    finally:
        await db.close()


ANALYSIS_LABELS = {
    "summary": "Resume",
    "key_points": "Points Cles",
    "actions": "Actions",
    "study_cards": "Fiches d'apprentissage",
    "quiz": "Quiz",
    "mindmap": "Carte mentale",
    "slides": "Slides",
    "infographic": "Infographie",
    "data_table": "Tableaux de donnees",
}


@router.get("/api/transcriptions/{tid}/export-analysis", dependencies=[Depends(verify_api_key)])
async def export_analysis(tid: int, type: str = Query(...)):
    """Telecharge une analyse individuelle en PDF."""
    from backend.outputs.pdf_export import generate_analysis_pdf

    db = await get_connection(config.DB_PATH)
    try:
        t = await get_transcription_with_segments(db, tid)
        if not t:
            raise HTTPException(404, "Transcription non trouvee.")
        analysis = await get_analysis(db, tid, type)
        if not analysis or not analysis.get("content"):
            raise HTTPException(404, f"Aucune analyse de type '{type}' trouvee.")
    finally:
        await db.close()

    label = ANALYSIS_LABELS.get(type, type)
    filename_base = t.get("filename", "audio").rsplit(".", 1)[0]
    pdf_bytes = generate_analysis_pdf(label, analysis["content"], t.get("filename", "audio"))

    return Response(
        pdf_bytes,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{filename_base}_{type}.pdf"'},
    )


@router.get("/api/transcriptions/{tid}/export-all-analyses", dependencies=[Depends(verify_api_key)])
async def export_all_analyses(tid: int):
    """Telecharge toutes les analyses en un seul PDF."""
    from backend.outputs.pdf_export import generate_all_analyses_pdf

    db = await get_connection(config.DB_PATH)
    try:
        t = await get_transcription_with_segments(db, tid)
        if not t:
            raise HTTPException(404, "Transcription non trouvee.")
        all_analyses = await get_analyses(db, tid)
    finally:
        await db.close()

    # Deduplicate: keep latest per type
    seen = {}
    for a in all_analyses:
        atype = a.get("type", "")
        if atype not in seen:
            seen[atype] = a

    if not seen:
        raise HTTPException(404, "Aucune analyse disponible.")

    # Build ordered dict label -> content
    analyses_content = {}
    for atype, label in ANALYSIS_LABELS.items():
        if atype in seen and seen[atype].get("content"):
            analyses_content[label] = seen[atype]["content"]

    filename_base = t.get("filename", "audio").rsplit(".", 1)[0]
    pdf_bytes = generate_all_analyses_pdf(t.get("filename", "audio"), analyses_content)

    return Response(
        pdf_bytes,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{filename_base}_analyses.pdf"'},
    )


@router.get("/api/transcriptions/{tid}/export-slides-pptx", dependencies=[Depends(verify_api_key)])
async def export_slides_pptx(tid: int):
    """Genere et telecharge les slides en PowerPoint (.pptx)."""
    db = await get_connection(config.DB_PATH)
    try:
        t = await get_transcription_with_segments(db, tid)
        if not t:
            raise HTTPException(404, "Transcription non trouvee.")
        analysis = await get_analysis(db, tid, "slides")
        if not analysis or not analysis.get("content"):
            raise HTTPException(404, "Aucune slide disponible. Generez les slides d'abord.")
    finally:
        await db.close()

    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    md = analysis["content"]
    slide_texts = [s.strip() for s in md.split("\n---\n") if s.strip()]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i, slide_md in enumerate(slide_texts):
        slide_lines = slide_md.split("\n")
        title = ""
        body_lines = []

        for line in slide_lines:
            stripped = line.strip()
            if stripped.startswith("# ") and not title:
                title = stripped[2:].strip()
            elif stripped.startswith("## "):
                title = stripped[3:].strip()
            else:
                body_lines.append(stripped)

        if i == 0:
            layout = prs.slide_layouts[0]  # Title slide
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title or "Presentation"
            subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
            if subtitle and body_lines:
                subtitle.text = "\n".join(l for l in body_lines if l)
        else:
            layout = prs.slide_layouts[1]  # Title + content
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title or f"Slide {i + 1}"
            content_ph = slide.placeholders[1] if len(slide.placeholders) > 1 else None
            if content_ph:
                tf = content_ph.text_frame
                tf.clear()
                first = True
                for line in body_lines:
                    if not line:
                        continue
                    text = line.lstrip("- ").replace("**", "")
                    if first:
                        tf.paragraphs[0].text = text
                        first = False
                    else:
                        p = tf.add_paragraph()
                        p.text = text
                    # Style
                    para = tf.paragraphs[-1]
                    para.font.size = Pt(18)
                    if line.startswith("- "):
                        para.level = 0

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    filename_base = t.get("filename", "audio").rsplit(".", 1)[0]
    return Response(
        buf.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename_base}_slides.pptx"'},
    )
